from flask import Flask, render_template, request, jsonify, session, make_response, current_app
from openai import OpenAI
import os
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
import markdown2
from docx import Document
import fitz  # PyMuPDF
import openpyxl
from pptx import Presentation
import io
from flask_migrate import Migrate
import uuid
from datetime import datetime, timedelta
from apscheduler.schedulers.background import BackgroundScheduler
from flask_session import Session
from redis import Redis
from rq import Queue
import json
import logging
import tiktoken

logging.basicConfig(level=logging.INFO)

app = Flask(__name__)
app.secret_key = 'assistant-ai-1a-urrugne-64122'  # Assurez-vous de garder votre secret key sécurisée et unique
app.config['SESSION_TYPE'] = 'redis'
app.config['SESSION_PERMANENT'] = False
app.config['SESSION_USE_SIGNER'] = True
app.config['SESSION_COOKIE_SECURE'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'None'

# Utiliser Redis pour les sessions et les tâches en file d'attente
redis_url = os.getenv('REDISCLOUD_URL')
if not redis_url:
    raise RuntimeError("REDIS_URL not set in the environment variables.")

redis_instance = Redis.from_url(redis_url)

# Tester la connexion à Redis test
try:
    redis_instance.ping()
    print("Connected to Redis successfully!")
except Exception as e:
    print(f"Failed to connect to Redis: {str(e)}")

app.config['SESSION_REDIS'] = redis_instance
Session(app)

q = Queue(connection=redis_instance)

CORS(app, supports_credentials=True, origins=['https://kokua.fr', 'https://www.kokua.fr'])

 

app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL').replace("://", "ql://", 1) # Remplacez pour corriger l'URL pour PostgreSQL
# Exemple de configuration temporaire pour la génération des migrations
# app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///temp.db'

app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)
migrate = Migrate(app, db)






class Conversation(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    session_id = db.Column(db.Text, unique=True, nullable=False)
    derniere_activite = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    messages = db.relationship('Message', backref='conversation', lazy=True)



class Message(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    conversation_id = db.Column(db.Integer, db.ForeignKey('conversation.id'), nullable=False)
    role = db.Column(db.String(50), nullable=False)  # "system", "user", ou "assistant"
    content = db.Column(db.Text, nullable=False)


client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))



def nettoyer_conversations_inactives():
    limite_inactivite = datetime.now() - timedelta(days=30)  # Exemple: 30 jours d'inactivité
    conversations_inactives = Conversation.query.filter(Conversation.derniere_activite < limite_inactivite).all()
    for conversation in conversations_inactives:
        db.session.delete(conversation)
    db.session.commit() 

scheduler = BackgroundScheduler()
scheduler.add_job(func=nettoyer_conversations_inactives, trigger="interval", days=1)  # Exécute tous les jours
scheduler.start()









def read_file_content(uploaded_file):
    file_type = uploaded_file.filename.split('.')[-1]
    content = ""
    if file_type == 'docx':
        document = Document(io.BytesIO(uploaded_file.read()))
        content = "\n".join([paragraph.text for paragraph in document.paragraphs])
    elif file_type == 'pdf':
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page in doc:
            content += page.get_text()
    elif file_type == 'txt':
        content = uploaded_file.read().decode('utf-8')
    elif file_type == 'xlsx':
        workbook = openpyxl.load_workbook(io.BytesIO(uploaded_file.read()))
        sheet = workbook.active
        data = []
        for row in sheet.iter_rows(values_only=True):
            data.append(' '.join([str(cell) for cell in row if cell is not None]))
        content = "\n".join(data)
    elif file_type == 'pptx':
        ppt = Presentation(io.BytesIO(uploaded_file.read()))
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        content = "\n".join(text)
    else:
        raise ValueError("Unsupported file type")

    return content, count_tokens(content)  # Return both content and its token count

    


@app.route('/')
def home(): 
    # session['message_history'] = []  # Réinitialise l'historique pour chaque nouvelle session
    return render_template('index.html')


# @app.route('/set-test-cookie') 
# def set_test_cookie():
#     response = make_response("Cookie Set")
#     response.set_cookie('test', 'chez jouni', secure=True, samesite='None', domain='.kokua.fr')
#     return response




@app.route('/reset-session', methods=['POST'])
def reset_session():
    session.pop('session_id', None)  # Supprime l'ID de session actuel
    return jsonify({"message": "Session réinitialisée"}), 200






# ! mise en place de rq test

@app.route('/ask', methods=['POST'])
def ask_question():
    session_id = session.get('session_id', str(uuid.uuid4()))
    session['session_id'] = session_id

    file_content, file_tokens = read_file_content(request.files['file']) if request.files.get('file') else (None, 0)

    data = {
        "config_key": request.form.get('config_key', 'default_config'),
        "question": request.form.get('question'),
        "session_id": session_id,
        "file_content": file_content,
        "file_tokens": file_tokens
    }
    app.logger.info(f"Requête reçue à /ask avec données : {data}")
    job = q.enqueue(process_ask_question, data)
    return jsonify({"job_id": job.get_id(), "session_id": session_id}), 202

def process_ask_question(data):
    with app.app_context():
        try:
            config_path = 'gpt_config.json'
            with open(config_path, 'r') as f:
                gpt_configs = json.load(f)
            app.logger.info(f"Configuration loaded for config_key: {data['config_key']}")

            gpt_config = gpt_configs.get(data['config_key'], {
                "model": "gpt-3.5-turbo",
                "temperature": 0.0,
                "max_tokens": "",
                "instructions": "Votre première réponse doit commencer par 'STAN :'",
                "top_p": 0.1,
                "frequency_penalty": 0,
                "presence_penalty": 0
            })
            app.logger.info(f"GPT Config: {gpt_config}")

            conversation = Conversation.query.filter_by(session_id=data['session_id']).first()
            if not conversation:
                conversation = Conversation(session_id=data['session_id'], derniere_activite=datetime.utcnow())
                db.session.add(conversation)
                db.session.commit()
                app.logger.info(f"New conversation created with session_id: {data['session_id']}")
            else:
                app.logger.info(f"Existing conversation found with session_id: {data['session_id']}")

            if data.get('file_content'):
                Message.query.filter_by(conversation_id=conversation.id).delete()
                db.session.commit()

            instructions_content = gpt_config['instructions']
            instructions_message = Message(conversation_id=conversation.id, role="system", content=instructions_content)
            db.session.add(instructions_message)
            db.session.commit()  # Commit after each add

            if data.get('question'):
                question_message = Message(conversation_id=conversation.id, role="user", content=data['question'])
                db.session.add(question_message)
                db.session.commit()  # Commit after each add

            if data.get('file_content'):
                file_content_message = Message(conversation_id=conversation.id, role="user", content=data['file_content'])
                db.session.add(file_content_message)
                db.session.commit()  # Commit after each add

            db_messages = Message.query.filter_by(conversation_id=conversation.id).all()
            messages_for_openai = [{"role": msg.role, "content": msg.content} for msg in db_messages]

            total_prompt_tokens = sum(count_tokens(msg['content'], gpt_config['model']) for msg in messages_for_openai) + data.get('file_tokens', 0)
            app.logger.info(f"Total prompt tokens calculated: {total_prompt_tokens}")

            max_tokens = gpt_config.get('max_tokens')
            if not max_tokens:
                max_tokens = calculate_max_tokens(total_prompt_tokens)
                app.logger.info(f"Dynamic max_tokens calculated: {max_tokens}")
            else:
                max_tokens = int(max_tokens)
                app.logger.info(f"Fixed max_tokens used: {max_tokens}")

            gpt_config['max_tokens'] = max_tokens

            response_html = handle_openai_request(gpt_config, messages_for_openai, conversation)
            app.logger.info(f"OpenAI response received")
            return {"response": response_html}
        except Exception as e:
            app.logger.error(f"Erreur lors du traitement de la requête : {e}")
            raise





# ! route indicateur qualité -----------------------------------------------



# ! CALCUL DES TOKENS AVEC TIKTOKEN -------------------

def count_tokens(text, model="gpt-4"):
    encoding = tiktoken.encoding_for_model(model)
    tokens = encoding.encode(text)
    return len(tokens)

def calculate_max_tokens(prompt_tokens, max_context_tokens=128000, max_output_tokens=4096):
    # Assure que la somme des tokens d'entrée et de sortie ne dépasse pas la fenêtre de contexte totale
    return min(max_context_tokens - prompt_tokens, max_output_tokens)




def calculate_quality_index(prompt_tokens, max_output_tokens, max_context_tokens=128000):
    available_output_tokens = min(max_context_tokens - prompt_tokens, max_output_tokens)
    ratio = prompt_tokens / available_output_tokens if available_output_tokens > 0 else float('inf')
    
    if ratio < 1:
        quality = "Élevée"
    elif 1 <= ratio < 5:
        quality = "Moyenne"
    elif 5 <= ratio < 10:
        quality = "Faible"
    else:
        quality = "Très faible"
    
    return quality, ratio


# ! FIN CALCUL DES TOKENS -------------------


@app.route('/quality_index', methods=['POST'])
def calculate_quality():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    file_content, file_tokens = read_file_content(file)
    
    max_output_tokens = 4096

    total_prompt_tokens = file_tokens  # Ajoutez d'autres tokens si nécessaire (instructions, etc.)
    
    quality, ratio = calculate_quality_index(total_prompt_tokens, max_output_tokens)
    
    quality_descriptions = {
       "Élevée": "La qualité du traitement est élevée. Le nombre de tokens d'entrée est bien dans les limites.",
       "Moyenne": "La qualité du traitement est moyenne. Le nombre de tokens d'entrée est proche des limites.",
       "Faible": "La qualité du traitement est faible. Le nombre de tokens d'entrée dépasse les limites recommandées.",
       "Très faible": "La qualité du traitement est très faible. Le nombre de tokens d'entrée dépasse largement les limites recommandées."
    }
    
    description = quality_descriptions.get(quality, "Qualité inconnue.")
    
    return jsonify({"quality": quality, "ratio": ratio, "description": description})



# ! fin route qualité ---------------------------------------------------------------------------






def process_messages(data, conversation):
    if data['question']:
        question_message = Message(conversation_id=conversation.id, role="user", content=data['question'])
        db.session.add(question_message)
    if data.get('file_content'):
        file_messages = [
            Message(conversation_id=conversation.id, role="user", content="Uploaded File: Provided"),
            Message(conversation_id=conversation.id, role="user", content=data['file_content'])
        ]
        db.session.extend(file_messages)
    db.session.commit()


def handle_openai_request(gpt_config, messages_for_openai, conversation):
    app.logger.info(f"Sending request to OpenAI with config: {gpt_config}")  # Log la configuration utilisée
    chat_completion = client.chat.completions.create(
        model=gpt_config['model'],
        messages=messages_for_openai,
        temperature=gpt_config['temperature'],
        max_tokens=gpt_config['max_tokens'],
        top_p=gpt_config['top_p'],
        frequency_penalty=gpt_config['frequency_penalty'],
        presence_penalty=gpt_config['presence_penalty'],
        
    )
    response_chatgpt = chat_completion.choices[0].message.content
    app.logger.info(f"OpenAI Response: {response_chatgpt}")  # Log la réponse d'OpenAI
    # Enregistrement de la réponse de l'assistant
    response_message = Message(conversation_id=conversation.id, role="assistant", content=response_chatgpt)
    db.session.add(response_message)
    db.session.commit()

    return markdown2.markdown(response_chatgpt)

@app.route('/results/<job_id>', methods=['GET'])
def get_results(job_id):
    job = q.fetch_job(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404

    if job.is_finished:
        result_data = job.result
        if isinstance(result_data, str):
            try:
                result_data = json.loads(result_data)
            except json.JSONDecodeError as e:
                return jsonify({"status": "error", "message": "Invalid JSON data", "details": str(e)}), 500

        response_content = result_data.get('response', "No response found")
        quality = result_data.get('quality', "Unknown")
        ratio = result_data.get('ratio', "N/A")

        return jsonify({"status": "finished", "response": response_content, "quality": quality, "ratio": ratio}), 200

    elif job.is_failed:
        return jsonify({"status": "failed", "error": "Job failed", "details": str(job.exc_info)}), 500

    else:
        return jsonify({"status": "processing"}), 202





@app.route('/conversation/<session_id>', methods=['GET'])
def get_conversation(session_id):
    conversation = Conversation.query.filter_by(session_id=session_id).first()
    if conversation:
        return jsonify({
            "session_id": conversation.session_id,
            "last_activity": conversation.derniere_activite,
            "messages": [{"id": msg.id, "content": msg.content, "role": msg.role} for msg in conversation.messages]
        }), 200
    else:
        return jsonify({"error": "Conversation not found"}), 404



@app.route('/ask_sync', methods=['POST'])
def ask_question_sync():
    session_id = session.get('session_id', str(uuid.uuid4()))
    session['session_id'] = session_id  # Assurez-vous que l'ID de session est bien conservé dans la session

    data = {
        "config_key": request.form.get('config_key', 'default_config'),  
        "question": request.form.get('question'),
        "session_id": session.get('session_id', str(uuid.uuid4())),
        "file_content": read_file_content(request.files.get('file')) if request.files.get('file') else None
    }
    app.logger.info(f"Requête reçue à /ask_sync avec données : {data}")

    response_data = process_ask_question(data)
    return jsonify({"response": response_data['response'], "session_id": session_id}), 200


# ! fonciton ping feu vert/rouge

@app.route('/ping', methods=['GET'])
def ping():
    return jsonify({"status": "ok"}), 200



if __name__ == '__main__':
    db.init_app(app)
    app.run(debug=True)
